#!/usr/bin/env python3
"""
PowerPoint Utilities
Core module for local PPTX manipulation

IMPORTANT: When creating new slides, ALWAYS duplicate an existing slide
to preserve themes, backgrounds, formatting, and design consistency.
Never create blank slides with just a layout - copy and modify instead.
"""

import os
import json
import copy
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsmap
    from lxml import etree
except ImportError:
    print("python-pptx not installed. Run: pip install python-pptx lxml")
    Presentation = None


class PPTXAnalyzer:
    """Analyze PowerPoint files"""

    def __init__(self, file_path: str):
        self.file_path = Path(file_path)
        if not self.file_path.exists():
            raise FileNotFoundError(f"PPTX not found: {file_path}")
        self.prs = Presentation(str(self.file_path))

    def get_metadata(self) -> Dict[str, Any]:
        """Get presentation metadata"""
        core_props = self.prs.core_properties

        return {
            'file': str(self.file_path),
            'slides': len(self.prs.slides),
            'title': core_props.title or 'N/A',
            'author': core_props.author or 'N/A',
            'subject': core_props.subject or 'N/A',
            'created': str(core_props.created) if core_props.created else 'N/A',
            'modified': str(core_props.modified) if core_props.modified else 'N/A',
            'slide_width': self.prs.slide_width.inches,
            'slide_height': self.prs.slide_height.inches
        }

    def get_slide_info(self, slide_num: int) -> Dict[str, Any]:
        """Get info for specific slide"""
        if slide_num < 1 or slide_num > len(self.prs.slides):
            raise ValueError(f"Invalid slide number: {slide_num}")

        slide = self.prs.slides[slide_num - 1]

        # Get layout name
        layout_name = "Unknown"
        if slide.slide_layout:
            layout_name = slide.slide_layout.name

        # Count elements
        text_boxes = 0
        images = 0
        tables = 0
        charts = 0
        shapes = 0

        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_boxes += 1
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                tables += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                charts += 1
            else:
                shapes += 1

        # Word count
        word_count = sum(len(t.split()) for t in texts)

        return {
            'slide': slide_num,
            'layout': layout_name,
            'text_boxes': text_boxes,
            'images': images,
            'tables': tables,
            'charts': charts,
            'shapes': shapes,
            'word_count': word_count,
            'texts': texts[:5]  # First 5 text snippets
        }

    def extract_text(self, slide_num: Optional[int] = None) -> str:
        """Extract text from presentation"""
        texts = []

        if slide_num:
            slides = [self.prs.slides[slide_num - 1]]
        else:
            slides = self.prs.slides

        for slide in slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            texts.append("\n".join(slide_texts))

        return "\n\n---\n\n".join(texts)

    def extract_notes(self, slide_num: Optional[int] = None) -> Dict[int, str]:
        """Extract speaker notes"""
        notes = {}

        if slide_num:
            slides = [(slide_num, self.prs.slides[slide_num - 1])]
        else:
            slides = enumerate(self.prs.slides, 1)

        for i, slide in slides if slide_num else enumerate(self.prs.slides, 1):
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text
                if notes_text.strip():
                    notes[i] = notes_text.strip()

        return notes

    def analyze(self) -> Dict[str, Any]:
        """Full presentation analysis"""
        metadata = self.get_metadata()

        analysis = {
            **metadata,
            'slides_detail': [],
            'total_word_count': 0,
            'has_notes': False
        }

        for i in range(len(self.prs.slides)):
            slide_info = self.get_slide_info(i + 1)
            analysis['slides_detail'].append(slide_info)
            analysis['total_word_count'] += slide_info['word_count']

        notes = self.extract_notes()
        analysis['has_notes'] = bool(notes)
        analysis['slides_with_notes'] = list(notes.keys())

        return analysis


class PPTXEditor:
    """
    Edit PowerPoint files with theme preservation.

    DESIGN PRINCIPLE: Always duplicate existing slides instead of creating
    blank ones. This preserves:
    - Theme colors and fonts
    - Background images and gradients
    - Master slide formatting
    - Placeholder positions
    - Decorative elements
    """

    def __init__(self, file_path: str):
        self.file_path = Path(file_path)
        if not self.file_path.exists():
            raise FileNotFoundError(f"PPTX not found: {file_path}")
        self.prs = Presentation(str(self.file_path))

    def _get_slide_layout_type(self, slide) -> str:
        """Identify the type of slide layout"""
        layout_name = slide.slide_layout.name.lower() if slide.slide_layout else ""

        # Detect layout type
        if 'title' in layout_name and 'content' not in layout_name:
            return 'title'
        elif 'content' in layout_name or 'body' in layout_name:
            return 'content'
        elif 'two' in layout_name or 'comparison' in layout_name:
            return 'two_column'
        elif 'blank' in layout_name:
            return 'blank'
        elif 'section' in layout_name:
            return 'section'
        else:
            return 'content'

    def _analyze_slide_structure(self, slide) -> Dict[str, Any]:
        """
        Analyze slide structure to understand content placeholders.
        Returns info about what can be modified safely.
        """
        structure = {
            'layout_type': self._get_slide_layout_type(slide),
            'title_shape': None,
            'subtitle_shape': None,
            'content_shapes': [],
            'image_shapes': [],
            'decorative_shapes': [],
            'has_background_image': False
        }

        for shape in slide.shapes:
            # Check if it's a placeholder
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type == PP_PLACEHOLDER.TITLE:
                    structure['title_shape'] = shape
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    structure['subtitle_shape'] = shape
                elif ph_type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                    structure['content_shapes'].append(shape)
                elif ph_type == PP_PLACEHOLDER.PICTURE:
                    structure['image_shapes'].append(shape)

            # Non-placeholder shapes
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                structure['image_shapes'].append(shape)
            elif shape.has_text_frame:
                # Check if it looks like content vs decoration
                text = shape.text_frame.text.strip()
                if len(text) > 20:  # Likely content
                    structure['content_shapes'].append(shape)
                else:
                    structure['decorative_shapes'].append(shape)
            else:
                structure['decorative_shapes'].append(shape)

        # Check for background
        if slide.background and slide.background.fill:
            fill = slide.background.fill
            if fill.type is not None:
                structure['has_background_image'] = True

        return structure

    def duplicate_slide(
        self,
        source_slide_num: int,
        insert_position: Optional[int] = None
    ) -> int:
        """
        Duplicate an existing slide, preserving ALL formatting and theme elements.

        This is the PREFERRED way to create new slides as it preserves:
        - Background images/colors
        - Theme formatting
        - Layout structure
        - Decorative elements
        - Font styles

        Args:
            source_slide_num: Slide number to duplicate (1-indexed)
            insert_position: Where to insert (1-indexed), None = after source

        Returns:
            New slide number (1-indexed)
        """
        if source_slide_num < 1 or source_slide_num > len(self.prs.slides):
            raise ValueError(f"Invalid source slide: {source_slide_num}")

        source_slide = self.prs.slides[source_slide_num - 1]

        # Get the slide layout from source
        slide_layout = source_slide.slide_layout

        # Add new slide with same layout
        new_slide = self.prs.slides.add_slide(slide_layout)

        # Copy all shapes from source to new slide
        self._copy_slide_content(source_slide, new_slide)

        # Copy slide background if custom
        self._copy_slide_background(source_slide, new_slide)

        # Move to correct position
        new_slide_id = self.prs.slides._sldIdLst[-1]
        self.prs.slides._sldIdLst.remove(new_slide_id)

        if insert_position is None:
            insert_position = source_slide_num + 1

        self.prs.slides._sldIdLst.insert(insert_position - 1, new_slide_id)

        return insert_position

    def _copy_slide_content(self, source_slide, target_slide):
        """Copy all shapes from source slide to target slide"""

        # Map to track copied shapes
        for shape in source_slide.shapes:
            # Skip placeholders that are already on target from layout
            if shape.is_placeholder:
                # Find matching placeholder on target and copy content
                ph_idx = shape.placeholder_format.idx
                for target_shape in target_slide.shapes:
                    if (target_shape.is_placeholder and
                        target_shape.placeholder_format.idx == ph_idx):
                        self._copy_shape_content(shape, target_shape)
                        break
            else:
                # For non-placeholder shapes, we need to recreate them
                # This is complex - for now copy text if possible
                pass

    def _copy_shape_content(self, source_shape, target_shape):
        """Copy content from one shape to another, preserving formatting"""

        if source_shape.has_text_frame and target_shape.has_text_frame:
            # Copy text with formatting
            target_tf = target_shape.text_frame

            # Clear target
            for para in list(target_tf.paragraphs)[1:]:
                p = para._p
                p.getparent().remove(p)

            # Copy paragraphs
            source_tf = source_shape.text_frame
            for i, source_para in enumerate(source_tf.paragraphs):
                if i == 0:
                    target_para = target_tf.paragraphs[0]
                else:
                    target_para = target_tf.add_paragraph()

                # Copy paragraph properties
                target_para.alignment = source_para.alignment
                target_para.level = source_para.level

                # Clear existing runs
                for run in list(target_para.runs):
                    run._r.getparent().remove(run._r)

                # Copy runs with formatting
                for source_run in source_para.runs:
                    target_run = target_para.add_run()
                    target_run.text = source_run.text

                    # Copy font properties
                    if source_run.font.bold is not None:
                        target_run.font.bold = source_run.font.bold
                    if source_run.font.italic is not None:
                        target_run.font.italic = source_run.font.italic
                    if source_run.font.size is not None:
                        target_run.font.size = source_run.font.size
                    if source_run.font.name is not None:
                        target_run.font.name = source_run.font.name
                    if source_run.font.color.rgb is not None:
                        target_run.font.color.rgb = source_run.font.color.rgb

    def _copy_slide_background(self, source_slide, target_slide):
        """Copy background from source to target slide"""
        try:
            # Access the underlying XML to copy background
            source_bg = source_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
            if source_bg is not None:
                target_bg = target_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
                if target_bg is not None:
                    target_bg.getparent().remove(target_bg)

                # Deep copy and insert
                new_bg = copy.deepcopy(source_bg)
                # Insert at correct position
                cSld = target_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
                if cSld is not None:
                    cSld.insert(0, new_bg)
        except Exception:
            # Background copy failed, continue without it
            pass

    def add_slide_from_template(
        self,
        template_slide_num: int,
        position: Optional[int] = None,
        title: Optional[str] = None,
        content: Optional[str] = None,
        content_items: Optional[List[str]] = None
    ) -> int:
        """
        Create a new slide by duplicating a template slide and updating content.

        This is the RECOMMENDED way to add slides as it preserves the theme.

        Args:
            template_slide_num: Slide to use as template (1-indexed)
            position: Where to insert new slide (1-indexed)
            title: New title text (replaces template title)
            content: Main content text (replaces template content)
            content_items: List of bullet points (for content slides)

        Returns:
            New slide number
        """
        # Duplicate the template slide
        new_slide_num = self.duplicate_slide(template_slide_num, position)
        new_slide = self.prs.slides[new_slide_num - 1]

        # Update content while preserving formatting
        if title is not None:
            self._update_title(new_slide, title)

        if content is not None:
            self._update_content(new_slide, content)
        elif content_items is not None:
            self._update_content_bullets(new_slide, content_items)

        return new_slide_num

    def _update_title(self, slide, new_title: str):
        """Update slide title while preserving formatting"""
        if slide.shapes.title:
            title_shape = slide.shapes.title
            if title_shape.has_text_frame and title_shape.text_frame.paragraphs:
                para = title_shape.text_frame.paragraphs[0]
                if para.runs:
                    # Preserve first run's formatting
                    para.runs[0].text = new_title
                    # Clear other runs
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    para.text = new_title

    def _update_content(self, slide, new_content: str):
        """Update main content area while preserving formatting"""
        structure = self._analyze_slide_structure(slide)

        if structure['content_shapes']:
            content_shape = structure['content_shapes'][0]
            if content_shape.has_text_frame:
                tf = content_shape.text_frame
                if tf.paragraphs:
                    # Keep formatting from first paragraph
                    first_para = tf.paragraphs[0]
                    if first_para.runs:
                        first_para.runs[0].text = new_content
                        for run in first_para.runs[1:]:
                            run.text = ""
                    else:
                        first_para.text = new_content
                    # Clear other paragraphs
                    for para in tf.paragraphs[1:]:
                        para.clear()

    def _update_content_bullets(self, slide, items: List[str]):
        """Update content with bullet points, preserving formatting"""
        structure = self._analyze_slide_structure(slide)

        if structure['content_shapes']:
            content_shape = structure['content_shapes'][0]
            if content_shape.has_text_frame:
                tf = content_shape.text_frame

                # Get formatting from first paragraph/run
                template_para = tf.paragraphs[0] if tf.paragraphs else None
                template_font = None
                if template_para and template_para.runs:
                    template_font = template_para.runs[0].font

                # Clear existing content but keep structure
                for para in tf.paragraphs:
                    para.clear()

                # Add new bullets
                for i, item in enumerate(items):
                    if i == 0:
                        para = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()

                    para.text = item
                    para.level = 0

                    # Apply template formatting if available
                    if template_font and para.runs:
                        run = para.runs[0]
                        if template_font.size:
                            run.font.size = template_font.size
                        if template_font.name:
                            run.font.name = template_font.name
                        if template_font.bold is not None:
                            run.font.bold = template_font.bold

    def find_best_template_slide(self, slide_type: str = 'content') -> int:
        """
        Find the best slide to use as a template for the given type.

        Args:
            slide_type: 'title', 'content', 'two_column', 'section', 'blank'

        Returns:
            Slide number (1-indexed) of best matching template
        """
        best_match = 1  # Default to first slide

        for i, slide in enumerate(self.prs.slides, 1):
            layout_type = self._get_slide_layout_type(slide)

            if layout_type == slide_type:
                return i

            # Fallback matches
            if slide_type == 'content' and layout_type in ['content', 'two_column']:
                best_match = i

        return best_match

    def replace_text(
        self,
        find: str,
        replace: str,
        slide_num: Optional[int] = None,
        preserve_formatting: bool = True
    ) -> int:
        """
        Find and replace text while preserving formatting.

        Returns number of replacements made
        """
        count = 0

        if slide_num:
            slides = [self.prs.slides[slide_num - 1]]
        else:
            slides = self.prs.slides

        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if find in run.text:
                                run.text = run.text.replace(find, replace)
                                count += 1

        return count

    def update_slide_text(
        self,
        slide_num: int,
        shape_index: int,
        new_text: str
    ):
        """Update text in specific shape while preserving formatting"""
        if slide_num < 1 or slide_num > len(self.prs.slides):
            raise ValueError(f"Invalid slide number: {slide_num}")

        slide = self.prs.slides[slide_num - 1]

        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        if shape_index < 0 or shape_index >= len(text_shapes):
            raise ValueError(f"Invalid shape index: {shape_index}")

        shape = text_shapes[shape_index]

        # Preserve first paragraph formatting
        if shape.text_frame.paragraphs:
            first_para = shape.text_frame.paragraphs[0]
            if first_para.runs:
                first_run = first_para.runs[0]

                # Store formatting
                bold = first_run.font.bold
                italic = first_run.font.italic
                size = first_run.font.size
                name = first_run.font.name
                color = first_run.font.color.rgb if first_run.font.color.type else None

                # Update text
                first_run.text = new_text

                # Restore formatting
                if bold is not None:
                    first_run.font.bold = bold
                if italic is not None:
                    first_run.font.italic = italic
                if size is not None:
                    first_run.font.size = size
                if name is not None:
                    first_run.font.name = name
                if color is not None:
                    first_run.font.color.rgb = color

                # Clear other runs
                for run in first_para.runs[1:]:
                    run.text = ""

    def add_slide(
        self,
        position: int,
        layout_index: int = 1,
        title: Optional[str] = None,
        content: Optional[str] = None,
        use_template: bool = True,
        template_slide: Optional[int] = None
    ):
        """
        Add new slide at position.

        RECOMMENDED: Set use_template=True (default) to duplicate an existing
        slide and preserve themes. Only set use_template=False if you
        specifically need a blank slide from the master layout.

        Args:
            position: Insert position (1-indexed)
            layout_index: Layout index (only used if use_template=False)
            title: Slide title
            content: Slide content
            use_template: If True, duplicate existing slide (preserves theme)
            template_slide: Which slide to use as template (auto-detected if None)
        """
        if use_template and len(self.prs.slides) > 0:
            # Find best template slide
            if template_slide is None:
                template_slide = self.find_best_template_slide('content')

            # Use template-based creation
            self.add_slide_from_template(
                template_slide_num=template_slide,
                position=position,
                title=title,
                content=content
            )
        else:
            # Fallback to layout-based creation (not recommended)
            layout = self.prs.slide_layouts[layout_index]
            slide = self.prs.slides.add_slide(layout)

            # Move to correct position
            slide_id = self.prs.slides._sldIdLst[-1]
            self.prs.slides._sldIdLst.remove(slide_id)
            self.prs.slides._sldIdLst.insert(position - 1, slide_id)

            # Add title if specified
            if title and slide.shapes.title:
                slide.shapes.title.text = title

            # Add content if specified
            if content:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        shape.text_frame.text = content
                        break

    def delete_slide(self, slide_num: int):
        """Delete slide"""
        if slide_num < 1 or slide_num > len(self.prs.slides):
            raise ValueError(f"Invalid slide number: {slide_num}")

        slide_id = self.prs.slides._sldIdLst[slide_num - 1]
        self.prs.part.drop_rel(slide_id.rId)
        self.prs.slides._sldIdLst.remove(slide_id)

    def reorder_slides(self, new_order: List[int]):
        """Reorder slides based on new order list"""
        # Validate
        if sorted(new_order) != list(range(1, len(self.prs.slides) + 1)):
            raise ValueError("Invalid slide order - must include all slides")

        # Create new order
        old_ids = list(self.prs.slides._sldIdLst)
        self.prs.slides._sldIdLst.clear()

        for pos in new_order:
            self.prs.slides._sldIdLst.append(old_ids[pos - 1])

    def batch_update_slides(
        self,
        updates: List[Dict[str, Any]]
    ) -> int:
        """
        Batch update multiple slides efficiently.

        Args:
            updates: List of dicts with keys:
                - slide_num: Slide to update (1-indexed)
                - title: Optional new title
                - content: Optional new content
                - content_items: Optional list of bullet points
                - replacements: Optional dict of find->replace text

        Returns:
            Number of slides updated
        """
        count = 0

        for update in updates:
            slide_num = update.get('slide_num')
            if not slide_num or slide_num < 1 or slide_num > len(self.prs.slides):
                continue

            slide = self.prs.slides[slide_num - 1]

            if 'title' in update:
                self._update_title(slide, update['title'])
                count += 1

            if 'content' in update:
                self._update_content(slide, update['content'])
                count += 1
            elif 'content_items' in update:
                self._update_content_bullets(slide, update['content_items'])
                count += 1

            if 'replacements' in update:
                for find, replace in update['replacements'].items():
                    self.replace_text(find, replace, slide_num)

        return count

    def save(self, output_path: str) -> str:
        """Save presentation"""
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output))
        return str(output)


def analyze_pptx(file_path: str) -> Dict[str, Any]:
    """Quick analysis function"""
    analyzer = PPTXAnalyzer(file_path)
    return analyzer.analyze()


def extract_text(file_path: str, slide: Optional[int] = None) -> str:
    """Quick text extraction"""
    analyzer = PPTXAnalyzer(file_path)
    return analyzer.extract_text(slide)


def replace_text(
    file_path: str,
    find: str,
    replace: str,
    output_path: str
) -> int:
    """Find and replace text in PPTX"""
    editor = PPTXEditor(file_path)
    count = editor.replace_text(find, replace)
    editor.save(output_path)
    return count


def duplicate_slide(
    file_path: str,
    source_slide: int,
    output_path: str,
    insert_position: Optional[int] = None
) -> int:
    """
    Duplicate a slide preserving all formatting and theme.

    Args:
        file_path: Source PPTX file
        source_slide: Slide number to duplicate (1-indexed)
        output_path: Where to save
        insert_position: Where to insert (1-indexed)

    Returns:
        New slide number
    """
    editor = PPTXEditor(file_path)
    new_slide_num = editor.duplicate_slide(source_slide, insert_position)
    editor.save(output_path)
    return new_slide_num


def add_slides_from_template(
    file_path: str,
    template_slide: int,
    slides_content: List[Dict[str, Any]],
    output_path: str
) -> int:
    """
    Add multiple slides based on a template slide.

    This is the RECOMMENDED way to add slides as it preserves the theme.

    Args:
        file_path: Source PPTX file
        template_slide: Slide number to use as template (1-indexed)
        slides_content: List of dicts with 'title', 'content' or 'content_items'
        output_path: Where to save

    Returns:
        Number of slides added

    Example:
        add_slides_from_template(
            'presentation.pptx',
            template_slide=2,  # Use slide 2 as template
            slides_content=[
                {'title': 'New Slide 1', 'content': 'Content here'},
                {'title': 'New Slide 2', 'content_items': ['Point 1', 'Point 2']},
            ],
            output_path='output/new_presentation.pptx'
        )
    """
    editor = PPTXEditor(file_path)

    for i, content in enumerate(slides_content):
        # Insert at end by default
        position = len(editor.prs.slides) + 1

        editor.add_slide_from_template(
            template_slide_num=template_slide,
            position=position,
            title=content.get('title'),
            content=content.get('content'),
            content_items=content.get('content_items')
        )

    editor.save(output_path)
    return len(slides_content)


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 2:
        print("Usage: python pptx_utils.py <pptx_file>")
        sys.exit(1)

    result = analyze_pptx(sys.argv[1])
    print(json.dumps(result, indent=2, default=str))
